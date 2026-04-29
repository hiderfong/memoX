"""文档解析器 - 支持 PDF、Markdown、TXT、DOCX"""

import asyncio
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Document:
    """文档对象"""
    id: str
    filename: str
    content: str
    metadata: dict = field(default_factory=dict)

    @property
    def extension(self) -> str:
        return Path(self.filename).suffix.lower()


@dataclass
class TextChunk:
    """文本块"""
    id: str
    content: str
    metadata: dict = field(default_factory=dict)
    index: int = 0


class BaseParser(ABC):
    """解析器基类"""

    @abstractmethod
    async def parse(self, file_path: Path, doc_id: str) -> Document:
        """解析文件"""
        pass

    @abstractmethod
    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        """分块"""
        pass


class MarkdownParser(BaseParser):
    """Markdown 解析器"""

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        content = file_path.read_text(encoding="utf-8")
        return Document(
            id=doc_id,
            filename=file_path.name,
            content=content,
            metadata={
                "type": "markdown",
                "path": str(file_path),
            }
        )

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        return self._text_chunk(document.content, document.id, document.metadata, chunk_size, overlap)

    def _text_chunk(self, text: str, doc_id: str, metadata: dict, chunk_size: int, overlap: int) -> list[TextChunk]:
        chunks: list[TextChunk] = []
        start = 0
        index = 0

        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end]

            # 尝试在句号或换行处分割
            if end < len(text):
                last_period = chunk_text.rfind("。")
                last_newline = chunk_text.rfind("\n")
                last_punctuation = max(last_period, last_newline)

                if last_punctuation > chunk_size * 0.5:
                    chunk_text = chunk_text[:last_punctuation + 1]
                    end = start + len(chunk_text)

            chunks.append(TextChunk(
                id=f"{doc_id}_chunk_{index}",
                content=chunk_text.strip(),
                metadata={**metadata, "chunk_index": index},
                index=index,
            ))

            index += 1
            start = end - overlap
            if start <= chunks[-1].index * (chunk_size - overlap):
                break

        return chunks


class TextParser(MarkdownParser):
    """纯文本解析器"""

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        content = file_path.read_text(encoding="utf-8")
        return Document(
            id=doc_id,
            filename=file_path.name,
            content=content,
            metadata={
                "type": "text",
                "path": str(file_path),
            }
        )


class PDFParser(BaseParser):
    """PDF 解析器"""

    # 解析限制
    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
    MAX_PAGES = 500  # 最大页数
    PARSE_TIMEOUT = 60  # 解析超时时间（秒）

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        try:
            import pymupdf
        except ImportError:
            raise ImportError("请安装 pymupdf: pip install pymupdf") from None

        # 检查文件大小
        file_size = file_path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"文件过大: {file_size / 1024 / 1024:.1f}MB > {self.MAX_FILE_SIZE / 1024 / 1024}MB")

        def _parse_sync():
            content_parts: list[str] = []
            metadata: dict = {}

            with pymupdf.open(file_path) as doc:
                page_count = len(doc)
                metadata = {
                    "pages": min(page_count, self.MAX_PAGES),
                    "title": doc.metadata.get("title", ""),
                }

                # 限制页数
                for page_num, page in enumerate(doc):
                    if page_num >= self.MAX_PAGES:
                        content_parts.append(f"[... 省略 {page_count - page_num} 页 ...]")
                        break

                    text = page.get_text()
                    if text.strip():
                        content_parts.append(f"[第 {page_num + 1} 页]\n{text}")

            return Document(
                id=doc_id,
                filename=file_path.name,
                content="\n\n".join(content_parts),
                metadata={**metadata, "type": "pdf", "path": str(file_path), "file_size": file_size},
            )

        # 在线程池中运行同步的 PDF 解析，带超时
        loop = asyncio.get_running_loop()
        try:
            return await asyncio.wait_for(
                loop.run_in_executor(None, _parse_sync),
                timeout=self.PARSE_TIMEOUT
            )
        except asyncio.TimeoutError:
            raise TimeoutError(f"PDF 解析超时（>{self.PARSE_TIMEOUT}秒），文件可能过于复杂") from None

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        return self._text_chunk(document.content, document.id, document.metadata, chunk_size, overlap)

    def _text_chunk(self, text: str, doc_id: str, metadata: dict, chunk_size: int, overlap: int) -> list[TextChunk]:
        # 按页分割
        pages = re.split(r'\[第 \d+ 页\]', text)
        chunks: list[TextChunk] = []
        current_chunk = ""
        index = 0

        for page_text in pages:
            if not page_text.strip():
                continue

            if len(current_chunk) + len(page_text) <= chunk_size:
                current_chunk += page_text + "\n\n"
            else:
                if current_chunk.strip():
                    chunks.append(TextChunk(
                        id=f"{doc_id}_chunk_{index}",
                        content=current_chunk.strip(),
                        metadata={**metadata, "chunk_index": index},
                        index=index,
                    ))
                    index += 1

                current_chunk = page_text + "\n\n"

                while len(current_chunk) > chunk_size:
                    chunks.append(TextChunk(
                        id=f"{doc_id}_chunk_{index}",
                        content=current_chunk[:chunk_size],
                        metadata={**metadata, "chunk_index": index},
                        index=index,
                    ))
                    index += 1
                    current_chunk = current_chunk[chunk_size - overlap:]

        if current_chunk.strip():
            chunks.append(TextChunk(
                id=f"{doc_id}_chunk_{index}",
                content=current_chunk.strip(),
                metadata={**metadata, "chunk_index": index},
                index=index,
            ))

        return chunks


class DOCXParser(BaseParser):
    """DOCX 解析器"""

    # 解析限制
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    MAX_PARAGRAPHS = 10000  # 最大段落数
    MAX_TABLES = 100  # 最大表格数
    MAX_TABLE_ROWS = 1000  # 单个表格最大行数
    MAX_TABLE_CELLS = 100  # 单行最大单元格数
    PARSE_TIMEOUT = 30  # 解析超时时间（秒）

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        try:
            import docx
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx") from None

        # 检查文件大小
        file_size = file_path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"文件过大: {file_size / 1024 / 1024:.1f}MB > {self.MAX_FILE_SIZE / 1024 / 1024}MB")

        def _parse_sync():
            doc = docx.Document(file_path)
            paragraphs: list[str] = []

            # 限制段落数
            for i, para in enumerate(doc.paragraphs):
                if i >= self.MAX_PARAGRAPHS:
                    paragraphs.append(f"[... 省略 {len(doc.paragraphs) - i} 个段落 ...]")
                    break
                text = para.text.strip()
                if text:
                    paragraphs.append(text)

            # 提取表格（带限制）
            tables_text: list[str] = []
            for table_idx, table in enumerate(doc.tables):
                if table_idx >= self.MAX_TABLES:
                    tables_text.append(f"[... 省略 {len(doc.tables) - table_idx} 个表格 ...]")
                    break

                for row_idx, row in enumerate(table.rows):
                    if row_idx >= self.MAX_TABLE_ROWS:
                        tables_text.append(f"[... 省略 {len(table.rows) - row_idx} 行 ...]")
                        break

                    # 限制单元格数
                    cells = []
                    for cell_idx, cell in enumerate(row.cells):
                        if cell_idx >= self.MAX_TABLE_CELLS:
                            cells.append(f"[+{len(row.cells) - cell_idx} more]")
                            break
                        cells.append(cell.text.strip())

                    if any(cells):
                        tables_text.append(" | ".join(cells))
                tables_text.append("")

            content = "\n\n".join(paragraphs)
            if tables_text:
                content += "\n\n### 表格内容 ###\n" + "\n".join(tables_text)

            return Document(
                id=doc_id,
                filename=file_path.name,
                content=content,
                metadata={
                    "type": "docx",
                    "paragraphs": len(paragraphs),
                    "tables": min(len(doc.tables), self.MAX_TABLES),
                    "path": str(file_path),
                    "file_size": file_size,
                }
            )

        # 在线程池中运行同步的 DOCX 解析，带超时
        loop = asyncio.get_running_loop()
        try:
            return await asyncio.wait_for(
                loop.run_in_executor(None, _parse_sync),
                timeout=self.PARSE_TIMEOUT
            )
        except asyncio.TimeoutError:
            raise TimeoutError(f"DOCX 解析超时（>{self.PARSE_TIMEOUT}秒），文件可能过于复杂") from None

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        return self._text_chunk(document.content, document.id, document.metadata, chunk_size, overlap)

    def _text_chunk(self, text: str, doc_id: str, metadata: dict, chunk_size: int, overlap: int) -> list[TextChunk]:
        chunks: list[TextChunk] = []
        start = 0
        index = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))

            # 在段落边界分割（仅当未到达末尾时）
            if end < len(text):
                last_newline = text.rfind("\n\n", start, end)
                if last_newline > start:
                    end = last_newline + 2

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(TextChunk(
                    id=f"{doc_id}_chunk_{index}",
                    content=chunk_text,
                    metadata={**metadata, "chunk_index": index},
                    index=index,
                ))
                index += 1

            # 到达文本末尾，退出循环，防止 start 回退导致无限循环
            if end >= len(text):
                break

            new_start = end - overlap
            # 确保 start 严格向前推进，防止 overlap 过大时死循环
            start = new_start if new_start > start else end

        return chunks


class WebPageParser(BaseParser):
    """网页 URL 解析器 - 抓取并提取正文文本"""

    MAX_CONTENT_SIZE = 5 * 1024 * 1024   # 响应体最大 5MB
    FETCH_TIMEOUT = 30                    # 请求超时秒数
    MAX_TEXT_LENGTH = 500 * 1024          # 提取文本上限 500KB

    # 抓取时忽略的标签（导航、广告、脚本等噪声）
    NOISE_TAGS = {
        "script", "style", "noscript", "nav", "header", "footer",
        "aside", "advertisement", "iframe", "svg", "form",
    }

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        """file_path.name 存放 URL，实际从网络抓取"""
        url = file_path.name  # 约定：调用方将 URL 写入 filename
        return await self.fetch_url(url, doc_id)

    async def fetch_url(self, url: str, doc_id: str) -> Document:
        """抓取 URL 并提取正文"""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("请安装 beautifulsoup4: pip install beautifulsoup4") from None

        import httpx

        headers = {
            "User-Agent": (
                "Mozilla/5.0 (compatible; MemoX/1.0; +https://github.com/memox)"
            ),
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
        }

        async with httpx.AsyncClient(
            follow_redirects=True,
            timeout=httpx.Timeout(self.FETCH_TIMEOUT),
            headers=headers,
        ) as client:
            response = await client.get(url)
            response.raise_for_status()

            content_type = response.headers.get("content-type", "")
            if len(response.content) > self.MAX_CONTENT_SIZE:
                raise ValueError(
                    f"页面内容过大: {len(response.content) / 1024 / 1024:.1f}MB"
                )

            # 非 HTML 内容直接返回文本
            if "html" not in content_type and "xml" not in content_type:
                text = response.text[: self.MAX_TEXT_LENGTH]
                return Document(
                    id=doc_id,
                    filename=url,
                    content=text,
                    metadata={"type": "webpage", "url": url, "content_type": content_type},
                )

            html = response.text

        # 解析 HTML
        soup = BeautifulSoup(html, "html.parser")

        # 提取元信息
        title = ""
        if soup.title and soup.title.string:
            title = soup.title.string.strip()

        description = ""
        meta_desc = soup.find("meta", attrs={"name": "description"})
        if meta_desc:
            description = meta_desc.get("content", "").strip()

        # 移除噪声标签
        for tag in soup.find_all(self.NOISE_TAGS):
            tag.decompose()

        # 优先提取语义化正文区域
        main_content = (
            soup.find("main")
            or soup.find("article")
            or soup.find(id=re.compile(r"content|main|article|post", re.I))
            or soup.find(class_=re.compile(r"content|main|article|post|body", re.I))
            or soup.find("body")
            or soup
        )

        # 提取文本，保留段落结构
        lines: list[str] = []
        if title:
            lines.append(f"# {title}")
        if description:
            lines.append(f"{description}\n")

        for element in main_content.descendants:
            if element.name in ("p", "h1", "h2", "h3", "h4", "h5", "li", "td", "th", "blockquote"):
                text = element.get_text(separator=" ", strip=True)
                if text and len(text) > 10:
                    lines.append(text)

        content = "\n\n".join(dict.fromkeys(lines))  # 去重并保序
        if len(content) > self.MAX_TEXT_LENGTH:
            content = content[: self.MAX_TEXT_LENGTH] + "\n[... 内容已截断 ...]"

        if not content.strip():
            # 兜底：获取全部文本
            content = soup.get_text(separator="\n", strip=True)[: self.MAX_TEXT_LENGTH]

        return Document(
            id=doc_id,
            filename=url,
            content=content,
            metadata={
                "type": "webpage",
                "url": url,
                "title": title,
                "description": description,
            },
        )

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        return self._text_chunk(document.content, document.id, document.metadata, chunk_size, overlap)

    def _text_chunk(self, text: str, doc_id: str, metadata: dict, chunk_size: int, overlap: int) -> list[TextChunk]:
        chunks: list[TextChunk] = []
        start = 0
        index = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))

            if end < len(text):
                last_newline = text.rfind("\n\n", start, end)
                if last_newline > start:
                    end = last_newline + 2

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(TextChunk(
                    id=f"{doc_id}_chunk_{index}",
                    content=chunk_text,
                    metadata={**metadata, "chunk_index": index},
                    index=index,
                ))
                index += 1

            if end >= len(text):
                break

            new_start = end - overlap
            start = new_start if new_start > start else end

        return chunks


class XLSXParser(BaseParser):
    """XLSX / XLS 解析器"""

    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    MAX_SHEETS = 20
    MAX_ROWS = 5000       # 每个 Sheet 最大行数
    MAX_COLS = 100        # 每行最大列数
    PARSE_TIMEOUT = 30

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl") from None

        file_size = file_path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"文件过大: {file_size / 1024 / 1024:.1f}MB > {self.MAX_FILE_SIZE / 1024 / 1024}MB")

        def _parse_sync():
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts: list[str] = []

            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                if sheet_idx >= self.MAX_SHEETS:
                    parts.append(f"[... 省略 {len(wb.sheetnames) - sheet_idx} 个 Sheet ...]")
                    break

                ws = wb[sheet_name]
                parts.append(f"### Sheet: {sheet_name} ###")
                row_count = 0

                for row in ws.iter_rows(values_only=True):
                    if row_count >= self.MAX_ROWS:
                        parts.append("[... 省略更多行 ...]")
                        break
                    cells = []
                    for col_idx, cell in enumerate(row):
                        if col_idx >= self.MAX_COLS:
                            cells.append(f"[+{len(row) - col_idx} more]")
                            break
                        if cell is not None:
                            cells.append(str(cell))
                    if any(c.strip() for c in cells):
                        parts.append(" | ".join(cells))
                        row_count += 1

            wb.close()
            return Document(
                id=doc_id,
                filename=file_path.name,
                content="\n".join(parts),
                metadata={
                    "type": "xlsx",
                    "sheets": min(len(wb.sheetnames), self.MAX_SHEETS),
                    "path": str(file_path),
                    "file_size": file_size,
                },
            )

        loop = asyncio.get_running_loop()
        try:
            return await asyncio.wait_for(
                loop.run_in_executor(None, _parse_sync),
                timeout=self.PARSE_TIMEOUT,
            )
        except asyncio.TimeoutError:
            raise TimeoutError(f"XLSX 解析超时（>{self.PARSE_TIMEOUT}秒）") from None

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        return self._text_chunk(document.content, document.id, document.metadata, chunk_size, overlap)

    def _text_chunk(self, text: str, doc_id: str, metadata: dict, chunk_size: int, overlap: int) -> list[TextChunk]:
        """按 Sheet 边界优先分块"""
        chunks: list[TextChunk] = []
        start = 0
        index = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))

            if end < len(text):
                # 优先在 Sheet 边界分割
                sheet_boundary = text.rfind("### Sheet:", start, end)
                if sheet_boundary > start:
                    end = sheet_boundary
                else:
                    last_newline = text.rfind("\n", start, end)
                    if last_newline > start:
                        end = last_newline + 1

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(TextChunk(
                    id=f"{doc_id}_chunk_{index}",
                    content=chunk_text,
                    metadata={**metadata, "chunk_index": index},
                    index=index,
                ))
                index += 1

            if end >= len(text):
                break

            new_start = end - overlap
            start = new_start if new_start > start else end

        return chunks


class PPTXParser(BaseParser):
    """PPTX 解析器"""

    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB（含图片）
    MAX_SLIDES = 200
    PARSE_TIMEOUT = 30

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx") from None

        file_size = file_path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"文件过大: {file_size / 1024 / 1024:.1f}MB > {self.MAX_FILE_SIZE / 1024 / 1024}MB")

        def _parse_sync():
            from pptx import Presentation

            prs = Presentation(file_path)
            parts: list[str] = []
            total_slides = len(prs.slides)

            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= self.MAX_SLIDES:
                    parts.append(f"[... 省略 {total_slides - slide_idx} 张幻灯片 ...]")
                    break

                slide_parts: list[str] = [f"[第 {slide_idx + 1} 页]"]

                for shape in slide.shapes:
                    # 文本框 / 标题 / 内容占位符
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_parts.append(text)
                    # 表格
                    elif shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            if any(cells):
                                slide_parts.append(" | ".join(cells))

                if len(slide_parts) > 1:  # 有实际内容才添加
                    parts.append("\n".join(slide_parts))

            return Document(
                id=doc_id,
                filename=file_path.name,
                content="\n\n".join(parts),
                metadata={
                    "type": "pptx",
                    "slides": min(total_slides, self.MAX_SLIDES),
                    "path": str(file_path),
                    "file_size": file_size,
                },
            )

        loop = asyncio.get_running_loop()
        try:
            return await asyncio.wait_for(
                loop.run_in_executor(None, _parse_sync),
                timeout=self.PARSE_TIMEOUT,
            )
        except asyncio.TimeoutError:
            raise TimeoutError(f"PPTX 解析超时（>{self.PARSE_TIMEOUT}秒）") from None

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        return self._text_chunk(document.content, document.id, document.metadata, chunk_size, overlap)

    def _text_chunk(self, text: str, doc_id: str, metadata: dict, chunk_size: int, overlap: int) -> list[TextChunk]:
        """按幻灯片边界优先分块"""
        chunks: list[TextChunk] = []
        start = 0
        index = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))

            if end < len(text):
                # 优先在幻灯片边界分割
                slide_boundary = text.rfind("[第 ", start, end)
                if slide_boundary > start:
                    end = slide_boundary
                else:
                    last_newline = text.rfind("\n\n", start, end)
                    if last_newline > start:
                        end = last_newline + 2

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(TextChunk(
                    id=f"{doc_id}_chunk_{index}",
                    content=chunk_text,
                    metadata={**metadata, "chunk_index": index},
                    index=index,
                ))
                index += 1

            if end >= len(text):
                break

            new_start = end - overlap
            start = new_start if new_start > start else end

        return chunks


class ImageParser(BaseParser):
    """图片 OCR 解析器 - Qwen VL 主路径 + pytesseract 兜底"""

    MAX_FILE_SIZE = 20 * 1024 * 1024  # 20MB
    OCR_TIMEOUT = 30

    def __init__(
        self,
        dashscope_api_key: str = "",
        dashscope_base_url: str = "https://dashscope.aliyuncs.com/compatible-mode/v1",
    ):
        self._api_key = dashscope_api_key
        self._base_url = dashscope_base_url

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        file_size = file_path.stat().st_size
        if file_size > self.MAX_FILE_SIZE:
            raise ValueError(f"图片过大: {file_size / 1024 / 1024:.1f}MB > {self.MAX_FILE_SIZE / 1024 / 1024}MB")

        # 主路径: Qwen VL
        if self._api_key:
            try:
                text = await asyncio.wait_for(
                    self._ocr_qwen_vl(file_path),
                    timeout=self.OCR_TIMEOUT,
                )
                return Document(
                    id=doc_id,
                    filename=file_path.name,
                    content=text,
                    metadata={"type": "image", "path": str(file_path), "file_size": file_size, "ocr_method": "qwen-vl"},
                )
            except Exception as e:
                import logging
                logging.getLogger(__name__).warning(f"Qwen VL OCR failed, falling back to pytesseract: {e}")

        # 兜底: pytesseract
        text = self._ocr_pytesseract(file_path)
        return Document(
            id=doc_id,
            filename=file_path.name,
            content=text,
            metadata={"type": "image", "path": str(file_path), "file_size": file_size, "ocr_method": "pytesseract"},
        )

    async def _ocr_qwen_vl(self, file_path: Path) -> str:
        """调用 Qwen VL 进行 OCR"""
        import base64

        import httpx

        image_data = file_path.read_bytes()
        ext = file_path.suffix.lstrip(".").lower()
        if ext == "jpg":
            ext = "jpeg"
        b64 = base64.b64encode(image_data).decode()

        payload = {
            "model": "qwen-vl-plus",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/{ext};base64,{b64}"}},
                        {"type": "text", "text": "请提取这张图片中的所有文字内容。如果图片中没有文字，请描述图片的主要内容。"},
                    ],
                }
            ],
            "max_tokens": 2048,
        }

        async with httpx.AsyncClient(timeout=httpx.Timeout(self.OCR_TIMEOUT)) as client:
            response = await client.post(
                f"{self._base_url}/chat/completions",
                headers={
                    "Authorization": f"Bearer {self._api_key}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
            response.raise_for_status()
            data = response.json()
            return data["choices"][0]["message"]["content"]

    def _ocr_pytesseract(self, file_path: Path) -> str:
        """本地 pytesseract OCR 兜底"""
        try:
            import pytesseract
            from PIL import Image
            image = Image.open(file_path)
            return pytesseract.image_to_string(image, lang="chi_sim+eng")
        except ImportError:
            return f"(图片文件: {file_path.name}，OCR 不可用 — 请安装 pytesseract 和 Pillow)"
        except Exception as e:
            return f"(图片 OCR 失败: {e})"

    async def chunk(self, document: Document, chunk_size: int = 500, overlap: int = 50) -> list[TextChunk]:
        """OCR 文本分块"""
        chunks: list[TextChunk] = []
        text = document.content
        start = 0
        index = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))
            if end < len(text):
                last_newline = text.rfind("\n", start, end)
                if last_newline > start:
                    end = last_newline + 1
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(TextChunk(
                    id=f"{document.id}_chunk_{index}",
                    content=chunk_text,
                    metadata={**document.metadata, "chunk_index": index},
                    index=index,
                ))
                index += 1
            if end >= len(text):
                break
            new_start = end - overlap
            start = new_start if new_start > start else end

        return chunks if chunks else [TextChunk(
            id=f"{document.id}_chunk_0",
            content=document.content,
            metadata={**document.metadata, "chunk_index": 0},
            index=0,
        )]


class DocumentParser:
    """文档解析器工厂"""

    def __init__(self, dashscope_api_key: str = "", dashscope_base_url: str = ""):
        image_parser = ImageParser(
            dashscope_api_key=dashscope_api_key,
            dashscope_base_url=dashscope_base_url or "https://dashscope.aliyuncs.com/compatible-mode/v1",
        )
        self._parsers: dict[str, BaseParser] = {
            ".md": MarkdownParser(),
            ".markdown": MarkdownParser(),
            ".txt": TextParser(),
            ".pdf": PDFParser(),
            ".docx": DOCXParser(),
            ".xlsx": XLSXParser(),
            ".xls": XLSXParser(),
            ".pptx": PPTXParser(),
            ".png": image_parser,
            ".jpg": image_parser,
            ".jpeg": image_parser,
            ".webp": image_parser,
        }

    def get_parser(self, filename: str) -> BaseParser:
        """获取对应的解析器"""
        ext = Path(filename).suffix.lower()

        # 尝试精确匹配
        if ext in self._parsers:
            return self._parsers[ext]

        # 尝试不带点的扩展名
        if f".{ext.lstrip('.')}" in self._parsers:
            return self._parsers[f".{ext.lstrip('.')}"]

        # 默认使用文本解析器
        return TextParser()

    async def parse(self, file_path: Path, doc_id: str) -> Document:
        """解析文档"""
        parser = self.get_parser(file_path.name)
        return await parser.parse(file_path, doc_id)

    async def parse_and_chunk(
        self,
        file_path: Path,
        doc_id: str,
        chunk_size: int = 500,
        overlap: int = 50
    ) -> tuple[Document, list[TextChunk]]:
        """解析并分块"""
        parser = self.get_parser(file_path.name)
        document = await parser.parse(file_path, doc_id)
        chunks = await parser.chunk(document, chunk_size, overlap)
        return document, chunks

    @property
    def supported_extensions(self) -> list[str]:
        """支持的扩展名"""
        return list(self._parsers.keys())
